import logging
import time
from pathlib import Path
from typing import Optional

from ..schemas import DocumentResult, DocumentMetadata, ParsingOptions, ParsingMethod
from .base import BaseParser

logger = logging.getLogger(__name__)

class TextParser(BaseParser):
    """Parses plain text and markdown files."""
    
    def parse(self, file_path: Path, options: ParsingOptions) -> DocumentResult:
        start_time = time.time()
        file_path = Path(file_path)
        content = ""
        error_msg = None
        status = "success"
        
        try:
            # Try UTF-8 first
            try:
                with open(file_path, "r", encoding="utf-8") as f:
                    content = f.read()
            except UnicodeDecodeError:
                # Fallback to GBK
                with open(file_path, "r", encoding="gbk") as f:
                    content = f.read()
                    
        except Exception as e:
            logger.error(f"Failed to parse text file {file_path}: {e}")
            status = "error"
            error_msg = str(e)
            
        metadata = DocumentMetadata(
            file_size=file_path.stat().st_size if file_path.exists() else 0,
            file_type=file_path.suffix.lower(),
            parsing_method=ParsingMethod.TEXT,
            parsing_time=time.time() - start_time
        )
        
        return DocumentResult(
            content=content,
            metadata=metadata,
            status=status,
            error_message=error_msg
        )

class DocxParser(BaseParser):
    """Parses Word (.docx) files."""
    
    def parse(self, file_path: Path, options: ParsingOptions) -> DocumentResult:
        start_time = time.time()
        file_path = Path(file_path)
        content = ""
        error_msg = None
        status = "success"
        
        try:
            from docx import Document
            doc = Document(file_path)
            content = "\n".join(para.text for para in doc.paragraphs)
        except ImportError:
            error_msg = "python-docx library not installed."
            status = "error"
            logger.error(error_msg)
        except Exception as e:
            logger.error(f"Failed to parse docx {file_path}: {e}")
            status = "error"
            error_msg = str(e)

        metadata = DocumentMetadata(
            file_size=file_path.stat().st_size if file_path.exists() else 0,
            file_type=".docx",
            parsing_method=ParsingMethod.UNSTRUCTURED,
            parsing_time=time.time() - start_time
        )
        
        return DocumentResult(
            content=content,
            metadata=metadata,
            status=status,
            error_message=error_msg
        )

class PptxParser(BaseParser):
    """Parses PowerPoint (.pptx) files."""
    
    def parse(self, file_path: Path, options: ParsingOptions) -> DocumentResult:
        start_time = time.time()
        file_path = Path(file_path)
        content = ""
        error_msg = None
        status = "success"
        
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            text_runs = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_runs.append(shape.text)
            content = "\n".join(text_runs)
        except ImportError:
            error_msg = "python-pptx library not installed."
            status = "error"
            logger.error(error_msg)
        except Exception as e:
            logger.error(f"Failed to parse pptx {file_path}: {e}")
            status = "error"
            error_msg = str(e)

        metadata = DocumentMetadata(
            file_size=file_path.stat().st_size if file_path.exists() else 0,
            file_type=".pptx",
            parsing_method=ParsingMethod.UNSTRUCTURED,
            parsing_time=time.time() - start_time
        )
        
        return DocumentResult(
            content=content,
            metadata=metadata,
            status=status,
            error_message=error_msg
        )

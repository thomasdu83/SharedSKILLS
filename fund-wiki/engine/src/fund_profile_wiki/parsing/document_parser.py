"""Parse raw diligence documents into text."""

from __future__ import annotations

import io
import os
import platform
import re
import shutil
import time
from dataclasses import dataclass, field
from pathlib import Path

import pymupdf
from PIL import Image
from docx import Document
from pptx import Presentation

try:
    import pytesseract
except ImportError:  # pragma: no cover - exercised by user environments
    pytesseract = None  # type: ignore[assignment]


DEFAULT_TESSERACT_CMD = r"C:\Program Files\Tesseract-OCR\tesseract.exe"


@dataclass(frozen=True)
class ParseResult:
    """Plain-text parse result with operational metadata."""

    text: str
    status: str = "success"
    method: str = "text"
    warnings: list[str] = field(default_factory=list)
    error: str = ""
    page_count: int = 0
    chars: int = 0
    quality_reason: str = ""


class DocumentParser:
    """Parse supported document types into plain text."""

    OCR_DPI = int(os.getenv("FPW_OCR_DPI", "300"))
    OCR_LANG = os.getenv("FPW_OCR_LANG", "chi_sim+eng")
    OCR_ENABLED = os.getenv("FPW_OCR_ENABLED", "1").strip().lower() not in {
        "0",
        "false",
        "no",
    }
    OCR_PAGE_TIMEOUT_SECONDS = int(os.getenv("FPW_OCR_PAGE_TIMEOUT_SECONDS", "180"))
    OCR_MIN_TOTAL_CHARS = int(os.getenv("FPW_OCR_MIN_TOTAL_CHARS", "200"))
    OCR_MIN_CHARS_PER_PAGE = int(os.getenv("FPW_OCR_MIN_CHARS_PER_PAGE", "50"))
    OCR_MIN_VALID_CHAR_RATIO = float(os.getenv("FPW_OCR_MIN_VALID_CHAR_RATIO", "0.30"))
    OCR_MAX_REPEATED_SHORT_LINE_RATIO = float(
        os.getenv("FPW_OCR_MAX_REPEATED_SHORT_LINE_RATIO", "0.28")
    )
    OCR_REPEATED_SHORT_LINE_MAX_LEN = int(
        os.getenv("FPW_OCR_REPEATED_SHORT_LINE_MAX_LEN", "24")
    )

    @staticmethod
    def parse(path: str | Path) -> str:
        """Backward-compatible parse API returning text only."""
        return DocumentParser.parse_with_metadata(path).text

    @staticmethod
    def is_encrypted_pdf(path: str | Path) -> bool:
        """Return True when a PDF requires a password before text extraction."""
        file_path = Path(path)
        if file_path.suffix.lower() != ".pdf":
            return False
        doc = None
        try:
            doc = pymupdf.open(file_path)
            return bool(getattr(doc, "needs_pass", False))
        except Exception:  # pylint: disable=broad-except
            return False
        finally:
            if doc is not None:
                doc.close()

    @staticmethod
    def parse_with_metadata(path: str | Path, progress=None) -> ParseResult:
        file_path = Path(path)
        ext = file_path.suffix.lower()
        start_time = time.time()
        if ext == ".pdf":
            return DocumentParser._parse_pdf(
                file_path, progress=progress, start_time=start_time
            )
        if ext == ".txt":
            return DocumentParser._text_result(
                file_path.read_text(encoding="utf-8", errors="ignore"), "text"
            )
        if ext == ".md":
            return DocumentParser._text_result(
                file_path.read_text(encoding="utf-8", errors="ignore"), "text"
            )
        if ext == ".docx":
            return DocumentParser._text_result(
                DocumentParser._parse_docx(file_path), "docx"
            )
        if ext == ".pptx":
            return DocumentParser._text_result(
                DocumentParser._parse_pptx(file_path), "pptx"
            )
        raise ValueError(f"Unsupported file type: {ext}")

    @staticmethod
    def _text_result(text: str, method: str) -> ParseResult:
        cleaned = clean_text(text)
        return ParseResult(text=cleaned, method=method, chars=len(cleaned))

    @staticmethod
    def _parse_pdf(
        path: Path, progress=None, start_time: float | None = None
    ) -> ParseResult:
        warnings: list[str] = []
        start = start_time or time.time()
        doc = None
        try:
            doc = pymupdf.open(path)
            page_count = doc.page_count
            text_layout = DocumentParser._extract_pdf_text_layout(doc)
            is_valid, reason = DocumentParser._is_text_valid(text_layout, page_count)
            cleaned_layout = clean_text(text_layout)
            if is_valid:
                return ParseResult(
                    text=cleaned_layout,
                    method="layout",
                    page_count=page_count,
                    chars=len(cleaned_layout),
                    quality_reason=reason,
                )

            message = f"PDF layout text quality is low; OCR fallback needed ({reason})"
            warnings.append(message)
            if progress:
                progress.warning("parse", f"{message}: file={path.name}")

            if not DocumentParser.OCR_ENABLED:
                warnings.append(
                    "OCR is disabled by FPW_OCR_ENABLED=0; keeping layout text."
                )
                return ParseResult(
                    text=cleaned_layout,
                    status="failed",
                    method="layout",
                    warnings=warnings,
                    page_count=page_count,
                    chars=len(cleaned_layout),
                    quality_reason=reason,
                )

            available, missing_reason = DocumentParser.ocr_available()
            if not available:
                warnings.append(f"OCR fallback skipped: {missing_reason}")
                if progress:
                    progress.warning(
                        "ocr", f"skip file={path.name} reason={missing_reason}"
                    )
                return ParseResult(
                    text=cleaned_layout,
                    status="failed",
                    method="layout",
                    warnings=warnings,
                    error=missing_reason,
                    page_count=page_count,
                    chars=len(cleaned_layout),
                    quality_reason=reason,
                )

            ocr_text = DocumentParser._parse_pdf_ocr(
                path, page_count=page_count, progress=progress
            )
            cleaned_ocr = clean_text(ocr_text)
            if len(cleaned_ocr.strip()) > len(cleaned_layout.strip()):
                return ParseResult(
                    text=cleaned_ocr,
                    method="ocr",
                    warnings=warnings,
                    page_count=page_count,
                    chars=len(cleaned_ocr),
                    quality_reason=reason,
                )

            warnings.append(
                "OCR produced no better text than layout extraction; keeping layout text."
            )
            if progress:
                progress.warning(
                    "ocr",
                    f"layout retained file={path.name} ocr_chars={len(cleaned_ocr)} layout_chars={len(cleaned_layout)}",
                )
            return ParseResult(
                text=cleaned_layout,
                status="partial_success",
                method="layout",
                warnings=warnings,
                page_count=page_count,
                chars=len(cleaned_layout),
                quality_reason=reason,
            )
        except Exception as exc:  # pylint: disable=broad-except
            if progress:
                progress.warning(
                    "parse",
                    f"PDF layout parse failed; OCR fallback needed file={path.name} error={exc}",
                )
            available, missing_reason = DocumentParser.ocr_available()
            if not available:
                return ParseResult(
                    text="",
                    status="failed",
                    method="pdf",
                    warnings=[f"OCR fallback skipped: {missing_reason}"],
                    error=str(exc),
                    quality_reason="layout_parse_failed",
                )
            try:
                ocr_text = DocumentParser._parse_pdf_ocr(
                    path, page_count=0, progress=progress
                )
                cleaned = clean_text(ocr_text)
                return ParseResult(
                    text=cleaned,
                    method="ocr",
                    warnings=[
                        f"Layout parse failed after {time.time() - start:.1f}s; OCR was used."
                    ],
                    chars=len(cleaned),
                    quality_reason="layout_parse_failed",
                )
            except Exception as ocr_exc:  # pylint: disable=broad-except
                return ParseResult(
                    text="",
                    status="failed",
                    method="ocr",
                    warnings=[f"OCR failed: {ocr_exc}"],
                    error=str(ocr_exc),
                    quality_reason="layout_parse_failed",
                )
        finally:
            if doc is not None:
                doc.close()

    @staticmethod
    def _parse_pdf_ocr(path: Path, page_count: int = 0, progress=None) -> str:
        if pytesseract is None:
            raise RuntimeError("pytesseract is not installed")
        DocumentParser._configure_tesseract()
        doc = pymupdf.open(path)
        chunks: list[str] = []
        total_pages = page_count or doc.page_count
        try:
            if progress:
                progress.info(
                    "ocr",
                    f"start file={path.name} pages={total_pages} dpi={DocumentParser.OCR_DPI}",
                )
            for page_num in range(doc.page_count):
                if progress:
                    progress.info(
                        "ocr",
                        f"page {page_num + 1}/{total_pages} start file={path.name}",
                    )
                page = doc.load_page(page_num)
                pix = page.get_pixmap(dpi=DocumentParser.OCR_DPI)
                image = Image.open(io.BytesIO(pix.tobytes("png")))
                try:
                    text = pytesseract.image_to_string(
                        image,
                        lang=DocumentParser.OCR_LANG,
                        timeout=DocumentParser.OCR_PAGE_TIMEOUT_SECONDS,
                    )
                except RuntimeError as exc:
                    text = f"[OCR page {page_num + 1} failed: {exc}]"
                    if progress:
                        progress.warning(
                            "ocr",
                            f"page {page_num + 1}/{total_pages} failed file={path.name} error={exc}",
                        )
                chunks.append(text)
                if progress:
                    progress.info(
                        "ocr",
                        f"page {page_num + 1}/{total_pages} done file={path.name} chars={len(text)}",
                    )
            if progress:
                progress.info(
                    "ocr",
                    f"done file={path.name} chars={sum(len(chunk) for chunk in chunks)}",
                )
            return "\n".join(chunks)
        finally:
            doc.close()

    @staticmethod
    def _extract_pdf_text_layout(doc) -> str:
        pages = [page.get_text("text") for page in doc]
        return "\n".join(pages)

    @staticmethod
    def _is_text_valid(text: str, page_count: int) -> tuple[bool, str]:
        if not text or not text.strip():
            return False, "empty text"

        total_chars = len(text.strip())
        if total_chars < DocumentParser.OCR_MIN_TOTAL_CHARS:
            return (
                False,
                f"total chars too low ({total_chars} < {DocumentParser.OCR_MIN_TOTAL_CHARS})",
            )

        chars_per_page = total_chars / max(page_count, 1)
        if chars_per_page < DocumentParser.OCR_MIN_CHARS_PER_PAGE:
            return (
                False,
                f"chars per page too low ({chars_per_page:.1f} < {DocumentParser.OCR_MIN_CHARS_PER_PAGE})",
            )

        valid_chars = 0
        for char in text:
            if (
                "\u4e00" <= char <= "\u9fff"
                or char.isalnum()
                or char in ",.;:!?()[]{}<>-+*/=_\"'%"
            ):
                valid_chars += 1
        valid_ratio = valid_chars / max(len(text), 1)
        if valid_ratio < DocumentParser.OCR_MIN_VALID_CHAR_RATIO:
            return (
                False,
                f"valid char ratio too low ({valid_ratio:.1%} < {DocumentParser.OCR_MIN_VALID_CHAR_RATIO:.0%})",
            )

        def norm(value: str) -> str:
            return re.sub(r"\s+", "", value)

        lines = text.splitlines()
        counts: dict[str, int] = {}
        for line in lines:
            normalized = norm(line)
            if 2 <= len(normalized) <= DocumentParser.OCR_REPEATED_SHORT_LINE_MAX_LEN:
                counts[normalized] = counts.get(normalized, 0) + 1

        repeated_lines = 0
        for line in lines:
            normalized = norm(line)
            if (
                2 <= len(normalized) <= DocumentParser.OCR_REPEATED_SHORT_LINE_MAX_LEN
                and counts.get(normalized, 0) >= 3
            ):
                repeated_lines += 1
        repeated_ratio = repeated_lines / max(len(lines), 1)
        if (
            repeated_ratio >= DocumentParser.OCR_MAX_REPEATED_SHORT_LINE_RATIO
            and chars_per_page < (DocumentParser.OCR_MIN_CHARS_PER_PAGE * 2)
        ):
            return False, f"repeated short lines too high ({repeated_ratio:.1%})"

        return True, "valid"

    @staticmethod
    def ocr_available() -> tuple[bool, str]:
        if pytesseract is None:
            return False, "Python package pytesseract is not installed."
        command = DocumentParser._tesseract_command()
        if command and not Path(command).exists():
            found = shutil.which("tesseract")
            if found:
                return True, ""
            return False, f"Tesseract executable not found: {command}"
        if not command and shutil.which("tesseract") is None:
            return False, "Tesseract executable is not on PATH."
        try:
            DocumentParser._configure_tesseract()
            langs = set(pytesseract.get_languages(config=""))
        except Exception as exc:  # pylint: disable=broad-except
            return False, f"Tesseract is not usable: {exc}"
        required_langs = [part for part in DocumentParser.OCR_LANG.split("+") if part]
        missing = [lang for lang in required_langs if lang not in langs]
        if missing:
            return False, f"Tesseract language data missing: {', '.join(missing)}"
        return True, ""

    @staticmethod
    def _configure_tesseract() -> None:
        if pytesseract is None:
            return
        command = DocumentParser._tesseract_command()
        if command:
            pytesseract.pytesseract.tesseract_cmd = command

    @staticmethod
    def _tesseract_command() -> str:
        explicit = os.getenv("FPW_TESSERACT_CMD") or os.getenv("TESSERACT_CMD")
        if explicit:
            return explicit
        if platform.system() == "Windows" and Path(DEFAULT_TESSERACT_CMD).exists():
            return DEFAULT_TESSERACT_CMD
        found = shutil.which("tesseract")
        return found or (
            DEFAULT_TESSERACT_CMD if platform.system() == "Windows" else ""
        )

    @staticmethod
    def _parse_docx(path: Path) -> str:
        doc = Document(path)
        parts = [p.text for p in doc.paragraphs if p.text.strip()]
        for table in doc.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                if any(cells):
                    parts.append(" | ".join(cells))
        return "\n".join(parts)

    @staticmethod
    def _parse_pptx(path: Path) -> str:
        prs = Presentation(path)
        parts: list[str] = []
        for idx, slide in enumerate(prs.slides, start=1):
            slide_parts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_parts.append(shape.text.strip())
            if slide_parts:
                parts.append(f"# Slide {idx}\n" + "\n".join(slide_parts))
        return "\n\n".join(parts)


def clean_text(text: str) -> str:
    """Small deterministic cleanup suitable for source notes."""
    text = text.replace("\x00", "")
    lines = [re.sub(r"\s+", " ", line).strip() for line in text.splitlines()]
    compact = []
    last = None
    for line in lines:
        if not line:
            continue
        if line == last:
            continue
        compact.append(line)
        last = line
    return "\n".join(compact)
